
from __future__ import annotations
"""
tg_bot/utils/file_processor.py
──────────────────────────────
Multi-format file reader: PDF, Excel, Word, ZIP, CSV, TXT, JSON, and more.

v29.0.0:
  • PPTX support
  • .env / .ini / .toml reading
  • Better ZIP handling (recursive)
  • More code file extensions
  • File metadata extraction
  • Improved error messages in Persian
"""

# NOTE: Consider using arki_project.utils.feature_registry for optional imports

import csv
import io
import json
import logging
import zipfile
from pathlib import Path

# ── TITANIUM v29.0 Integration ──


logger = logging.getLogger(__name__)

# Maximum text length to return (avoid flooding AI context).
MAX_TEXT_LENGTH = 12000  # v7: increased from 8000

# Code/text extensions that can be read as plain text
_TEXT_EXTENSIONS = {
    ".txt", ".csv", ".json", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".html", ".css", ".scss", ".sass", ".less",
    ".md", ".rst", ".xml", ".yml", ".yaml", ".toml", ".ini", ".cfg",
    ".log", ".env", ".sh", ".bash", ".zsh", ".bat", ".ps1",
    ".sql", ".graphql", ".proto",
    ".java", ".kt", ".kts", ".scala", ".groovy",
    ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".swift",
    ".rb", ".php", ".pl", ".lua", ".r", ".m", ".mm",
    ".dockerfile", ".gitignore", ".editorconfig",
    ".vue", ".svelte",
}


async def process_file(file_bytes: bytes, filename: str) -> str:
    """
    Extract text content from a file based on its extension.

    Returns extracted text or an error message.
    """
    ext = Path(filename).suffix.lower()

    # Special case: files without extension
    basename = Path(filename).name.lower()
    if basename in ("dockerfile", "makefile", "gemfile", "rakefile", ".gitignore", ".env"):
        ext = ".txt"

    processors = {
        ".pdf": _read_pdf,
        ".xlsx": _read_excel,
        ".xls": _read_excel,
        ".docx": _read_docx,
        ".pptx": _read_pptx,
        ".zip": _read_zip,
        ".csv": _read_csv,
        ".json": _read_json,
    }

    # Check text extensions
    if ext in _TEXT_EXTENSIONS:
        processor = _read_text
    else:
        processor = processors.get(ext)

    if processor is None:
        supported = sorted(set(list(processors.keys()) + list(_TEXT_EXTENSIONS)))
        return (
            f"❌ فرمت پشتیبانی نمی‌شود: `{ext}`\n\n"
            f"فرمت‌های پشتیبانی شده:\n{', '.join(supported[:30])}"
        )

    try:
        text = processor(file_bytes, filename)
        file_size = len(file_bytes)
        size_str = _format_size(file_size)

        # Add file info header
        header = f"📄 *{filename}* ({size_str})\n{'─' * 30}\n\n"

        if len(text) > MAX_TEXT_LENGTH:
            text = text[:MAX_TEXT_LENGTH] + (
                f"\n\n… [ادامه حذف شده — {len(text):,} کاراکتر کل]"
            )
        return header + text

    except ImportError as e:
        package = _suggest_package(ext)
        return (
            f"❌ کتابخانه مورد نیاز برای `{ext}` نصب نیست: {e}\n"
            f"نصب: `pip install {package}`"
        )
    except Exception as e:
        logger.error("Error reading %s: %s", filename, e, exc_info=True)
        return f"❌ خطا در خواندن `{filename}`: {e}"


def _read_text(data: bytes, filename: str) -> str:
    """Read plain text files with multi-encoding support."""
    for encoding in ("utf-8", "utf-8-sig", "latin-1", "cp1256", "cp1252", "windows-1256"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


def _read_csv(data: bytes, filename: str) -> str:
    """Read CSV and format as readable table."""
    text = _read_text(data, filename)
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return "فایل CSV خالی است."

    # Format as markdown table.
    header = rows[0]
    lines = [" | ".join(header)]
    lines.append(" | ".join("---" for _ in header))
    for row in rows[1:50]:  # Limit to 50 rows.
        lines.append(" | ".join(row))

    result = f"📊 {len(rows) - 1} ردیف × {len(header)} ستون\n\n" + "\n".join(lines)
    if len(rows) > 51:
        result += f"\n\n… ({len(rows) - 51} ردیف دیگر)"
    return result


def _read_json(data: bytes, filename: str) -> str:
    """Read and pretty-print JSON."""
    text = _read_text(data, filename)
    parsed = json.loads(text)

    # Add structure info
    info = ""
    if isinstance(parsed, list):
        info = f"📋 آرایه JSON: {len(parsed)} آیتم\n\n"
    elif isinstance(parsed, dict):
        info = f"📋 آبجکت JSON: {len(parsed)} کلید\n\n"

    return info + json.dumps(parsed, indent=2, ensure_ascii=False)[:MAX_TEXT_LENGTH]


def _read_pdf(data: bytes, filename: str) -> str:
    """Read PDF using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        # Fallback: try pdfplumber.
        try:
            import pdfplumber

            with pdfplumber.open(io.BytesIO(data)) as pdf:
                total_pages = len(pdf.pages)
                pages = []
                for i, page in enumerate(pdf.pages[:30]):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"--- صفحه {i + 1} ---\n{text}")
                if not pages:
                    return "PDF بدون متن قابل استخراج."
                header = f"📕 PDF: {total_pages} صفحه\n\n"
                return header + "\n\n".join(pages)
        except ImportError:
            raise ImportError("PyMuPDF or pdfplumber")

    doc = fitz.open(stream=data, filetype="pdf")
    total_pages = len(doc)
    pages = []
    for i, page in enumerate(doc[:30]):
        text = page.get_text()
        if text.strip():
            pages.append(f"--- صفحه {i + 1} ---\n{text}")
    doc.close()

    if not pages:
        return "PDF بدون متن قابل استخراج."

    header = f"📕 PDF: {total_pages} صفحه\n\n"
    return header + "\n\n".join(pages)


def _read_excel(data: bytes, filename: str) -> str:
    """Read Excel file using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("openpyxl")

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets_text = []

    for sheet_name in wb.sheetnames[:5]:
        ws = wb[sheet_name]
        rows = []
        row_count = 0
        for row in ws.iter_rows(max_row=50, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(" | ".join(cells))
            row_count += 1

        if rows:
            sheets_text.append(
                f"📊 شیت: {sheet_name} ({row_count} ردیف)\n" + "\n".join(rows)
            )

    wb.close()

    if not sheets_text:
        return "فایل اکسل خالی است."

    return "\n\n".join(sheets_text)


def _read_docx(data: bytes, filename: str) -> str:
    """Read Word document using python-docx."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx")

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    if not paragraphs:
        return "فایل Word خالی است."

    return f"📝 Word: {len(paragraphs)} پاراگراف\n\n" + "\n\n".join(paragraphs)


def _read_pptx(data: bytes, filename: str) -> str:
    """Read PowerPoint presentation using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx")

    prs = Presentation(io.BytesIO(data))
    slides_text = []

    for i, slide in enumerate(prs.slides[:30]):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides_text.append(f"--- اسلاید {i + 1} ---\n" + "\n".join(texts))

    if not slides_text:
        return "فایل PowerPoint بدون متن."

    return f"📊 PPTX: {len(prs.slides)} اسلاید\n\n" + "\n\n".join(slides_text)


def _read_zip(data: bytes, filename: str) -> str:
    """List contents of a ZIP file and read small text files inside."""
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        file_list = zf.namelist()
        total_size = sum(zf.getinfo(n).file_size for n in file_list)
        result = (
            f"📦 ZIP: {len(file_list)} فایل "
            f"({_format_size(total_size)} حجم کل)\n\n"
        )
        result += "\n".join(f"  • {f}" for f in file_list[:50])
        if len(file_list) > 50:
            result += f"\n  … و {len(file_list) - 50} فایل دیگر"

        # Read small text files.
        read_count = 0
        for name in file_list:
            if read_count >= 5:
                break
            ext = "." + name.rsplit(".", 1)[-1].lower() if "." in name else ""
            if ext in _TEXT_EXTENSIONS:
                info = zf.getinfo(name)
                if info.file_size < 10000:
                    try:
                        content = zf.read(name).decode("utf-8", errors="replace")
                        result += f"\n\n{'─' * 30}\n📄 {name}:\n{content}"
                        read_count += 1
                    except Exception as e:
                        logger.debug("Suppressed: %s", e)

        return result


def _format_size(size_bytes: int) -> str:
    """Format file size to human-readable string."""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    if size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    return f"{size_bytes / (1024 * 1024):.1f} MB"


def _suggest_package(ext: str) -> str:
    """Suggest pip package for a file extension."""
    suggestions = {
        ".pdf": "PyMuPDF",
        ".xlsx": "openpyxl",
        ".xls": "openpyxl",
        ".docx": "python-docx",
        ".pptx": "python-pptx",
    }
    return suggestions.get(ext, "unknown")


